"""Generate Word (.docx) and PowerPoint (.pptx) versions of the Friday demo brief."""
from __future__ import annotations

from pathlib import Path

from docx import Document
from docx.shared import Pt, RGBColor, Inches
from pptx import Presentation
from pptx.util import Inches as PInches, Pt as PPt
from pptx.dml.color import RGBColor as PRGB
from pptx.enum.shapes import MSO_SHAPE

ROOT = Path(__file__).resolve().parents[2]
OUT = ROOT / "docs" / "demos"
OUT.mkdir(parents=True, exist_ok=True)

BRAND = RGBColor(0x0F, 0x3D, 0x7A)
PBRAND = PRGB(0x0F, 0x3D, 0x7A)


# =============================================================================
# Word document
# =============================================================================
def build_docx() -> Path:
    doc = Document()

    # Base styles
    style = doc.styles["Normal"]
    style.font.name = "Calibri"
    style.font.size = Pt(11)

    def h(level: int, text: str) -> None:
        p = doc.add_heading(text, level=level)
        for run in p.runs:
            run.font.color.rgb = BRAND

    def p(text: str) -> None:
        doc.add_paragraph(text)

    def b(text: str) -> None:
        doc.add_paragraph(text, style="List Bullet")

    def table(headers: list[str], rows: list[list[str]]) -> None:
        t = doc.add_table(rows=1 + len(rows), cols=len(headers))
        t.style = "Light Grid Accent 1"
        for i, hd in enumerate(headers):
            cell = t.rows[0].cells[i]
            cell.text = hd
            for run in cell.paragraphs[0].runs:
                run.bold = True
        for r, row in enumerate(rows, start=1):
            for c, val in enumerate(row):
                t.rows[r].cells[c].text = val
        doc.add_paragraph()

    # Cover
    h(0, "EGP Window — Friday Demo Brief")
    p("Audience: M42 (customer) — technical + architectural")
    p("Duration: 30–45 min   |   Presenter: Vijay   |   Date: Friday 2026-07-17")
    p("Repo: https://github.com/nandyap/EPG-MAF   |   Tag: w11-cutover (ee6204d)")

    # Elevator pitch
    h(1, "0. Elevator Pitch (30 sec)")
    p(
        "We have ported the customer's LangGraph clinical-genomics prototype to "
        "Microsoft Agent Framework (MAF) on Azure. All 11 workstreams are complete "
        "(421 unit tests passing), the HTTP surface is browsable locally via Swagger UI "
        "today, and the system is ready for a Compass key + Azure subscription to move "
        "from stub to live. Six open questions await customer answers — captured in "
        "epg-maf/docs/blockers.md."
    )

    # Agenda
    h(1, "1. Agenda")
    table(
        ["#", "Time", "Section"],
        [
            ["1", "5 min", "The problem & why MAF"],
            ["2", "5 min", "Architecture overview"],
            ["3", "10 min", "Live demo (Swagger UI)"],
            ["4", "10 min", "MAF components we use & why"],
            ["5", "5 min", "Parallel vs. sequential dispatch toggle"],
            ["6", "5 min", "What's shipped vs. what's blocked"],
            ["7", "5 min", "What we need from M42 (Q&A)"],
        ],
    )

    # Problem
    h(1, "2. The Problem (5 min)")
    p(
        "A clinical genomics decision-support agent. A clinician (or patient) asks a "
        "natural-language question about a single patient's genomic + clinical data. "
        "The agent must:"
    )
    b("Route the question to the right specialist(s): PRS, Genomic Variants, Family History, PGx, Phenotype.")
    b("Call the relevant tools against a curated clinical database (Postgres in prod, DuckDB seed for tests).")
    b("Synthesise a clinically defensible answer with structured provenance.")
    b("Enforce strict single-patient scope + PHI redaction + auditable trail.")
    p(
        "Prototype today: LangGraph — single-process, no auth, no observability budget, "
        "no resilience/retry, no clinical audit sink, no production topology."
    )
    p(
        "What we built: the same workflow reimplemented on MAF with Azure-native "
        "infrastructure, defence-in-depth, and full enterprise observability."
    )

    # Architecture
    h(1, "3. Architecture Overview (5 min)")
    h(2, "3.1 Runtime topology (production)")
    p(
        "Client → Azure Front Door (WAF) → APIM (retry + circuit-breaker) → "
        "Container Apps (FastAPI) → auth → workflow_runtime.run → ChatWorkflow → "
        "OrchestrationWorkflow → specialists → Postgres/Cosmos/Compass/Foundry."
    )
    h(2, "3.2 Chat vs. orchestration sub-workflow")
    b("ChatWorkflow: decides whether the question needs clinical data. If not, synthesise reply directly.")
    b("OrchestrationWorkflow: bounded loop (orch_router → dispatcher → specialists → joiner → orch_router) with iteration budget = 12.")

    h(2, "3.3 What each specialist does")
    table(
        ["Specialist", "Domain", "Example tool", "Result schema"],
        [
            ["PRS", "Polygenic risk scores", "get_patient_prs, search_prs_annotations", "PRSResultList"],
            ["Genomic Variants", "Variant calls + annotations", "get_patient_variants, search_variant_annotations", "VariantResultList"],
            ["Family History", "Pedigree criteria", "get_patient_family_history (privacy-redacted)", "FamilyHistoryResultList"],
            ["PGx", "Drug–gene interactions (CPIC)", "get_patient_pgx, search_pgx_annotations", "PGXResultList"],
            ["Phenotype", "HPO/MONDO phenotype match", "get_patient_phenotype", "PhenotypeResultList"],
        ],
    )

    # MAF components
    h(1, "4. MAF Components We Use, and Why (10 min)")

    for title, para in [
        (
            "4.1 Executor + WorkflowContext",
            "Every node is a class(Executor) with a single @handler method. Handlers receive typed messages "
            "and either send the next message or yield the terminal output. Deterministic, testable, replayable.",
        ),
        (
            "4.2 Structured decision types (Pydantic, extra='forbid')",
            "ChatRouterDecision, SpecialistDispatchSet, ChatRequestBody, ChatResponseBody. "
            "Every hand-off is validated at the boundary. Unknown fields are startup errors, not runtime ones.",
        ),
        (
            "4.3 DI Container",
            "build_container(settings) wires every singleton (DB pool, Cosmos, LLM clients, repositories, "
            "authenticator, workflow runtime). Enables no-key smoke server, unit tests, prod — all share shape.",
        ),
        (
            "4.4 Observability (OpenTelemetry — W08)",
            "7 span kinds; 10 KPI metrics; PHI-safe attribute allowlist enforced on every set_attribute() call "
            "plus a CI grep test to block regressions.",
        ),
        (
            "4.5 Resilience (W09)",
            "Typed error hierarchy → HTTP codes. RetryingSpecialistLlm decorator with exp+jitter, "
            "independent of APIM retry. Specialist isolation: one failure marks a slot as failed; "
            "the rest still return.",
        ),
        (
            "4.6 Auth + audit (W07)",
            "Prod: Entra JWT verification against JWKS + role check. Dev: StubAuthenticator that accepts a "
            "JSON payload (this is what makes today's demo possible without an Entra tenant). Every request "
            "emits an AuditEvent (route, clinician_id, patient_id, decision, outcome).",
        ),
    ]:
        h(2, title)
        p(para)

    # Dispatch toggle
    h(1, "5. Parallel vs. Sequential Dispatch (5 min)")
    h(2, "5.1 The toggle")
    b("Setting: ORCH_DISPATCH_MODE = sequential | parallel (config/settings.py)")
    b("Companion: ORCH_MAX_FANOUT_WIDTH — caps concurrent specialists")
    b("Read by: OrchRouterExecutor.handle_state")

    h(2, "5.2 What the router does")
    b("Sequential (default): one specialist per iteration. Safe, deterministic, cost-linear, easy to audit.")
    b("Parallel: a set of independent specialists per iteration. Faster wall-clock, higher peak cost.")
    b("If parallel decision comes in while mode is sequential → downgraded to first specialist with warning log.")

    h(2, "5.3 How to toggle after deploy")
    table(
        ["Environment", "Mechanism"],
        [
            ["Local dev", "Set ORCH_DISPATCH_MODE=parallel in .env, restart process"],
            ["Preprod / prod", "Edit infra/env/prod.bicepparam → redeploy; or change Container App env var → restart revision"],
            ["Runtime toggle (no restart)", "Not built — would need Azure App Configuration + feature flags"],
        ],
    )

    h(2, "5.4 Analytics partitioning")
    p(
        "Every orch_router.dispatched log record carries orch.mode and orch.width. W08 lifts them to OTEL "
        "span attributes so dashboards partition latency/cost/error-rate by mode without a config join."
    )

    # Shipped vs blocked
    h(1, "6. Shipped vs. Blocked (5 min)")
    h(2, "6.1 Shipped (11/11 workstreams, tagged w11-cutover)")
    table(
        ["WS", "Title", "Highlights"],
        [
            ["W01", "Foundation", "Settings, DI, PG pool, Cosmos, structured logs"],
            ["W02", "Clinical data layer", "Schema, seed, repositories, authz"],
            ["W03", "Domain repositories", "5 repos + 14 tool shims + family-history privacy"],
            ["W04", "MAF workflow skeleton", "Chat + orchestration sub-workflows, routing, budget"],
            ["W05", "Specialist agents", "5 specialists with tools + provenance"],
            ["W06", "Parallel exec & mode-parity", "Sequential ↔ parallel harness + parity tests"],
            ["W07", "Auth + authorization", "Entra JWT, stub for dev, AuditEvent"],
            ["W08", "Observability", "OTEL SDK, 7 spans, 10 metrics, PHI allowlist"],
            ["W09", "Resilience & errors", "Typed errors, retry, isolation, response formatter"],
            ["W10", "Testing, eval & load", "Golden-set schema, scorers, PHI CI, load runbook"],
            ["W11", "Cutover + runbooks", "FastAPI, APIM, 9 alerts, 3 dashboards, 8 runbooks, cutover playbook"],
        ],
    )
    p("Testing snapshot: 421 unit tests passing, 21 integration tests skipped (require live Postgres).")

    h(2, "6.2 Open blockers (need M42 answers)")
    table(
        ["ID", "Blocker", "Owner", "Blocks"],
        [
            ["B-001", "PRS 'EGP-evaluated' metadata model", "BIX + M42", "Golden items S4, R5; PRS disclosure"],
            ["B-002", "Identity & session model", "M42 + IAM", "Session-pinning; refusal wording"],
            ["B-003", "Patient identifier formats", "BIX + M42", "ScopeGuard regex coverage"],
            ["B-004", "Approved refusal message wording", "M42 UX + Clinical", "Golden-set assertion strings"],
            ["B-005", "Session lifecycle + logout contract", "M42 + Frontend", "TTL sizing; refusal usefulness"],
            ["B-006", "Audit sink + alert threshold", "M42 Security + SIEM", "Sev-3 alert rule; retention"],
        ],
    )

    h(2, "6.3 Golden-dataset gap analysis")
    b("Gap 1: Single-patient scope guardrail (G1–G5). Design ready. Blocked on B-002 + B-003 + B-004.")
    b("Gap 2: Annotation vs. patient-scan (G6–G9). Prompt rule + repository lint. No external blocker.")
    b("Gap 3: 'Not EGP-evaluated' PRS disclosure (S4, R5). Blocked on B-001.")

    # Live demo
    h(1, "7. Live Demo Script (10 min)")
    p("Setup (before demo): cd epg-maf && .\\.venv\\Scripts\\python.exe scripts\\serve_smoke.py")
    b("Health probe: browse http://127.0.0.1:8000/healthz → 200 OK")
    b("Swagger UI: http://127.0.0.1:8000/docs → click Authorize (top-right)")
    b('Paste token: {"oid":"demo","tid":"demo","roles":["Clinician"],"exp":9999999999}')
    b('POST /chat → body: {"thread_id":"T-demo","patient_id":"P001","message":"What PRS does this patient have?"}')
    b("Response: populated prs + pgx slots + synthesised reply")
    b("Point out response schema: trace_id, agents_completed, per-specialist slot with status/output/errors")
    b("Show terminal logs: chat_router.decided, orch_router.dispatched, specialist_joiner.merged, orch_router.terminal")
    b("If time permits: pytest -m 'not integration and not parity and not chaos' -q → 421 passed")

    # QA
    h(1, "8. Likely Q&A")
    table(
        ["Question", "Short answer"],
        [
            ["Why MAF over LangGraph?", "Azure-native, typed decisions, testable, first-class observability, Foundry-ready"],
            ["Cost profile?", "Bounded by iteration budget (12) + fanout width; sequential = cost-linear"],
            ["PHI risk?", "Static allowlist + CI grep; family-history results explicitly redacted"],
            ["What if a specialist fails?", "Slot marked failed; other specialists still run; reply still returns"],
            ["How do we roll back?", "cutover.md playbook — idempotent, 30s target via revision swap"],
            ["Real LLM demo when?", "Immediately after Compass key + Azure — one env var flip"],
            ["What is NOT built?", "Real prod deploy, Foundry judge, 30–50 BIX-approved golden items, Grafana"],
        ],
    )

    # Close
    h(1, "9. Post-demo Close (2 min)")
    b("Deliver blocker list; ask for a single point-of-contact per blocker on M42 side.")
    b("Ask for Compass subscription key + Azure subscription onboarding target date.")
    b("Ask when a BIX reviewer can commit to the 43-item golden set for real evaluation.")
    b("Schedule follow-up walk-through of the golden-set gaps once B-001 to B-006 are unblocked.")

    path = OUT / "friday-demo-brief.docx"
    doc.save(path)
    return path


# =============================================================================
# PowerPoint
# =============================================================================
def build_pptx() -> Path:
    prs = Presentation()
    prs.slide_width = PInches(13.333)
    prs.slide_height = PInches(7.5)

    BLANK = prs.slide_layouts[6]
    TITLE = prs.slide_layouts[0]

    def add_title_slide(title: str, subtitle: str) -> None:
        s = prs.slides.add_slide(TITLE)
        s.shapes.title.text = title
        s.placeholders[1].text = subtitle
        for r in s.shapes.title.text_frame.paragraphs[0].runs:
            r.font.color.rgb = PBRAND
            r.font.bold = True

    def add_content_slide(title: str, bullets: list[str], footer: str = "") -> None:
        s = prs.slides.add_slide(BLANK)
        # Title bar
        title_box = s.shapes.add_textbox(PInches(0.5), PInches(0.3), PInches(12.3), PInches(0.7))
        tf = title_box.text_frame
        tf.text = title
        p0 = tf.paragraphs[0]
        p0.runs[0].font.size = PPt(30)
        p0.runs[0].font.bold = True
        p0.runs[0].font.color.rgb = PBRAND

        # Bullets
        body = s.shapes.add_textbox(PInches(0.6), PInches(1.2), PInches(12.1), PInches(5.5))
        btf = body.text_frame
        btf.word_wrap = True
        for i, item in enumerate(bullets):
            para = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
            para.text = f"• {item}"
            para.font.size = PPt(18)
            para.space_after = PPt(6)

        if footer:
            fbox = s.shapes.add_textbox(PInches(0.5), PInches(6.9), PInches(12.3), PInches(0.4))
            ftf = fbox.text_frame
            ftf.text = footer
            ftf.paragraphs[0].runs[0].font.size = PPt(10)
            ftf.paragraphs[0].runs[0].font.italic = True
            ftf.paragraphs[0].runs[0].font.color.rgb = PBRAND

    def add_table_slide(title: str, headers: list[str], rows: list[list[str]]) -> None:
        s = prs.slides.add_slide(BLANK)
        title_box = s.shapes.add_textbox(PInches(0.5), PInches(0.3), PInches(12.3), PInches(0.7))
        title_box.text_frame.text = title
        p0 = title_box.text_frame.paragraphs[0]
        p0.runs[0].font.size = PPt(30)
        p0.runs[0].font.bold = True
        p0.runs[0].font.color.rgb = PBRAND

        table = s.shapes.add_table(
            rows=1 + len(rows),
            cols=len(headers),
            left=PInches(0.5),
            top=PInches(1.2),
            width=PInches(12.3),
            height=PInches(5.5),
        ).table
        for i, hd in enumerate(headers):
            c = table.cell(0, i)
            c.text = hd
            for r in c.text_frame.paragraphs[0].runs:
                r.font.bold = True
                r.font.size = PPt(14)
                r.font.color.rgb = PRGB(0xFF, 0xFF, 0xFF)
            c.fill.solid()
            c.fill.fore_color.rgb = PBRAND
        for r, row in enumerate(rows, start=1):
            for c, val in enumerate(row):
                cell = table.cell(r, c)
                cell.text = val
                for run in cell.text_frame.paragraphs[0].runs:
                    run.font.size = PPt(12)

    # 1 Cover
    add_title_slide(
        "EGP Window — Customer Demo",
        "Microsoft Agent Framework port  |  Friday 2026-07-17\n"
        "Tag: w11-cutover  •  421 tests passing  •  11 workstreams shipped",
    )

    # 2 Elevator pitch
    add_content_slide(
        "Elevator Pitch",
        [
            "LangGraph prototype ported to Microsoft Agent Framework on Azure.",
            "11/11 workstreams complete; 421 unit tests passing.",
            "HTTP surface browsable locally via Swagger UI today (no LLM key needed).",
            "Ready for Compass key + Azure subscription to move stub → live.",
            "6 open questions await M42 answers (see epg-maf/docs/blockers.md).",
        ],
    )

    # 3 Agenda
    add_table_slide(
        "Agenda (30–45 min)",
        ["#", "Time", "Section"],
        [
            ["1", "5 min", "The problem & why MAF"],
            ["2", "5 min", "Architecture overview"],
            ["3", "10 min", "Live demo (Swagger UI)"],
            ["4", "10 min", "MAF components we use & why"],
            ["5", "5 min", "Parallel vs. sequential dispatch toggle"],
            ["6", "5 min", "Shipped vs. blocked"],
            ["7", "5 min", "What we need from M42 (Q&A)"],
        ],
    )

    # 4 Problem
    add_content_slide(
        "The Problem",
        [
            "Clinical genomics decision-support agent — one patient per session.",
            "Route to the right specialist(s): PRS, Variants, Family History, PGx, Phenotype.",
            "Call tools against curated clinical DB (Postgres prod / DuckDB seed).",
            "Synthesise clinically defensible answer with structured provenance.",
            "Enforce single-patient scope, PHI redaction, auditable trail.",
            "Prototype: LangGraph — no auth, no observability, no resilience, no prod topology.",
        ],
    )

    # 5 Architecture
    add_content_slide(
        "Architecture — Runtime Topology",
        [
            "Client → Azure Front Door (WAF)",
            "→ APIM (retry + circuit-breaker)",
            "→ Container Apps (FastAPI create_app)",
            "→ auth → workflow_runtime.run",
            "   ChatWorkflow → OrchestrationWorkflow (loop, budget=12)",
            "→ 5 specialists → repositories",
            "Data plane: Postgres Flexible + Cosmos DB + Key Vault + Compass + Foundry",
            "Telemetry: App Insights + Log Analytics (7 span kinds, 10 KPI metrics)",
        ],
    )

    # 6 Specialists
    add_table_slide(
        "Specialists",
        ["Specialist", "Domain", "Result schema"],
        [
            ["PRS", "Polygenic risk scores", "PRSResultList"],
            ["Genomic Variants", "Variant calls + annotations", "VariantResultList"],
            ["Family History", "NCCN / Amsterdam II criteria (redacted)", "FamilyHistoryResultList"],
            ["PGx", "Drug–gene interactions (CPIC)", "PGXResultList"],
            ["Phenotype", "HPO / MONDO matching", "PhenotypeResultList"],
        ],
    )

    # 7 MAF components
    add_content_slide(
        "MAF Components & Why",
        [
            "Executor + WorkflowContext — every node is typed, replayable, unit-testable.",
            "Pydantic decisions (extra='forbid') — every hand-off validated at boundary.",
            "DI Container — one factory wires prod + tests + smoke server identically.",
            "OpenTelemetry — 7 span kinds, 10 metrics, PHI-safe attribute allowlist.",
            "Resilience — typed errors, RetryingSpecialistLlm, specialist isolation.",
            "Auth — Entra JWT (prod) / StubAuthenticator (dev), AuditEvent per request.",
        ],
    )

    # 8 Dispatch toggle
    add_content_slide(
        "Parallel vs. Sequential Dispatch",
        [
            "Setting: ORCH_DISPATCH_MODE = sequential | parallel",
            "Companion: ORCH_MAX_FANOUT_WIDTH — caps concurrent specialists",
            "Sequential (default): one specialist per iteration; safe, cost-linear.",
            "Parallel: independent set per iteration; faster wall-clock, higher peak cost.",
            "Router downgrades parallel decisions in sequential mode with a warning.",
            "Toggle in prod: change bicepparam → redeploy, OR Container App env var + restart.",
            "Every log line carries orch.mode + orch.width for dashboard partitioning.",
        ],
    )

    # 9 Shipped
    add_table_slide(
        "Shipped: 11/11 Workstreams",
        ["WS", "Title", "Highlights"],
        [
            ["W01", "Foundation", "DI, PG pool, Cosmos, structured logs"],
            ["W02", "Clinical data", "Schema, seed, repos, authz"],
            ["W03", "Repositories", "5 repos, 14 tools, FH privacy"],
            ["W04", "Workflow", "Chat + orchestration, routing, budget"],
            ["W05", "Specialists", "5 agents + provenance"],
            ["W06", "Mode parity", "Sequential ↔ parallel parity tests"],
            ["W07", "Auth", "Entra JWT + AuditEvent"],
            ["W08", "Observability", "OTEL, spans, metrics, PHI allowlist"],
            ["W09", "Resilience", "Errors + retry + isolation"],
            ["W10", "Testing", "Golden set, scorers, PHI CI"],
            ["W11", "Cutover", "APIM, alerts, dashboards, runbooks"],
        ],
    )

    # 10 Blockers
    add_table_slide(
        "Open Blockers (need M42 answers)",
        ["ID", "Blocker", "Blocks"],
        [
            ["B-001", "PRS 'EGP-evaluated' metadata", "Gap 3 disclosure logic"],
            ["B-002", "Identity & session model", "Session-pinning + refusal copy"],
            ["B-003", "Patient identifier formats", "ScopeGuard regex coverage"],
            ["B-004", "Approved refusal wording", "Golden-set assertion strings"],
            ["B-005", "Session lifecycle + logout", "TTL sizing"],
            ["B-006", "Audit sink + alert threshold", "Sev-3 alert rule"],
        ],
    )

    # 11 Golden gaps
    add_content_slide(
        "Golden-Dataset Gap Analysis",
        [
            "Gap 1: Single-patient scope guardrail (G1–G5) — design ready, blocked on B-002/B-003/B-004.",
            "Gap 2: Annotation vs. patient-scan (G6–G9) — prompt rule + repo lint test; no external blocker.",
            "Gap 3: 'Not EGP-evaluated' PRS disclosure (S4, R5) — blocked on B-001 schema question.",
            "Golden set has 43 items (BIX-drafted); real BIX review still pending.",
        ],
    )

    # 12 Live demo script
    add_content_slide(
        "Live Demo — Swagger UI",
        [
            "Setup: cd epg-maf && .\\.venv\\Scripts\\python.exe scripts\\serve_smoke.py",
            "GET  http://127.0.0.1:8000/healthz  → 200 OK",
            "Open http://127.0.0.1:8000/docs  → click Authorize",
            'Token: {"oid":"demo","tid":"demo","roles":["Clinician"],"exp":9999999999}',
            'Body: {"thread_id":"T-demo","patient_id":"P001","message":"What PRS…?"}',
            "→ populated prs + pgx slots + synthesised reply.",
            "Show logs: chat_router.decided → orch_router.dispatched → joiner → terminal → synthesise.",
        ],
    )

    # 13 Q&A prep
    add_table_slide(
        "Likely Q&A",
        ["Question", "Short answer"],
        [
            ["Why MAF over LangGraph?", "Azure-native, typed, testable, observable, Foundry-ready"],
            ["Cost profile?", "Iteration budget=12 + fanout width; sequential is cost-linear"],
            ["PHI risk?", "Attribute allowlist + CI grep + FH redaction"],
            ["Specialist failure?", "Isolated slot; other specialists + reply still return"],
            ["Rollback?", "cutover.md — idempotent, 30s target via revision swap"],
            ["Real LLM demo?", "Immediately after Compass key + Azure subscription"],
            ["What is NOT built?", "Prod deploy; Foundry judge; BIX-approved 30–50 golden items"],
        ],
    )

    # 14 Close
    add_content_slide(
        "Post-Demo Close & Asks",
        [
            "Deliver blocker list — ask for a single point-of-contact per blocker.",
            "Compass subscription key — when?",
            "Azure subscription onboarding target date?",
            "BIX reviewer commitment to the 43-item golden set — when?",
            "Schedule follow-up walk-through once B-001 to B-006 are unblocked.",
        ],
        footer="See docs/demos/friday-demo-brief.md for the full narrative.",
    )

    # 15 Thanks
    add_title_slide(
        "Thank You",
        "Questions?\n\nRepo: https://github.com/nandyap/EPG-MAF   |   Tag: w11-cutover",
    )

    path = OUT / "friday-demo-brief.pptx"
    prs.save(path)
    return path


if __name__ == "__main__":
    d = build_docx()
    p = build_pptx()
    print(f"Wrote {d}")
    print(f"Wrote {p}")
